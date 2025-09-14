"""
Document Parser with Gemini Embeddings and ChromaDB Storage
===========================================================

This module handles document parsing, chunking, embedding generation using Gemini,
and storage in ChromaDB vector database for RAG pipeline.
"""

import os
import re
import uuid
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pptx import Presentation
import PyPDF2
import logging

# Import our custom modules
from gemini_embeddings import GeminiEmbeddingService, EmbeddingResult
from chroma_vector_db import ChromaVectorDB, DocumentChunk

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class EnhancedDocumentParser:
    """
    Enhanced document parser with Gemini embeddings and ChromaDB storage
    Handles the complete RAG pipeline: Parse → Chunk → Embed → Store
    """
    
    def __init__(self, 
                 embedding_service: Optional[GeminiEmbeddingService] = None,
                 vector_db: Optional[ChromaVectorDB] = None):
        """
        Initialize the enhanced document parser
        
        Args:
            embedding_service: Gemini embedding service instance
            vector_db: ChromaDB vector database instance
        """
        self.supported_formats = ['.pdf', '.pptx', '.docx']
        
        # Initialize services
        self.embedding_service = embedding_service or GeminiEmbeddingService()
        self.vector_db = vector_db or ChromaVectorDB()
        
        logger.info("✅ Enhanced Document Parser initialized")
    
    def extract_text_from_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PDF with enhanced metadata
        
        Args:
            file_path: Path to PDF file
        
        Returns:
            List of page data with text and metadata
        """
        try:
            pages_data = []
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    
                    if text and text.strip():
                        cleaned_text = self._clean_text(text)
                        if cleaned_text:
                            pages_data.append({
                                'text': cleaned_text,
                                'page_number': page_num,
                                'type': 'page',
                                'word_count': len(cleaned_text.split()),
                                'char_count': len(cleaned_text)
                            })
            
            logger.info(f"📄 Extracted text from {len(pages_data)} PDF pages")
            return pages_data
            
        except Exception as e:
            logger.error(f"❌ Error extracting PDF text: {e}")
            raise
    
    def extract_text_from_ppt(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PowerPoint with enhanced metadata
        
        Args:
            file_path: Path to PPTX file
        
        Returns:
            List of slide data with text and metadata
        """
        try:
            presentation = Presentation(file_path)
            slides_data = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_texts = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    combined_text = " ".join(slide_texts)
                    cleaned_text = self._clean_text(combined_text)
                    
                    if cleaned_text:
                        slides_data.append({
                            'text': cleaned_text,
                            'slide_number': slide_num,
                            'type': 'slide',
                            'word_count': len(cleaned_text.split()),
                            'char_count': len(cleaned_text),
                            'shape_count': len(slide_texts)
                        })
            
            logger.info(f"🎞️ Extracted text from {len(slides_data)} PowerPoint slides")
            return slides_data
            
        except Exception as e:
            logger.error(f"❌ Error extracting PPT text: {e}")
            raise
    
    def _clean_text(self, text: str) -> str:
        """
        Enhanced text cleaning with better preprocessing
        
        Args:
            text: Raw text to clean
        
        Returns:
            Cleaned text
        """
        # Remove slide/page numbers
        text = re.sub(r'\b(slide|page)\s*\d+\b', '', text, flags=re.IGNORECASE)
        
        # Remove headers/footers
        text = re.sub(r'^(header|footer).*$', '', text, flags=re.MULTILINE | re.IGNORECASE)
        
        # Normalize whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive punctuation
        text = re.sub(r'[.]{3,}', '...', text)
        text = re.sub(r'[-]{3,}', '---', text)
        
        # Clean up bullet points
        text = re.sub(r'^\s*[•▪▫‣⁃]\s*', '', text, flags=re.MULTILINE)
        
        # Remove very short lines (likely noise)
        lines = [line.strip() for line in text.split('\n') if len(line.strip()) > 5]
        
        return ' '.join(lines).strip()
    
    def create_intelligent_chunks(self, pages_data: List[Dict[str, Any]], 
                                 source_file: str,
                                 chunk_size: int = 400,
                                 overlap: int = 50) -> List[DocumentChunk]:
        """
        Create intelligent text chunks with overlap and metadata
        
        Args:
            pages_data: List of page/slide data
            source_file: Source file name
            chunk_size: Target chunk size in words
            overlap: Overlap between chunks in words
        
        Returns:
            List of DocumentChunk objects
        """
        chunks = []
        
        for page_data in pages_data:
            text = page_data['text']
            words = text.split()
            
            # If text is shorter than chunk size, create single chunk
            if len(words) <= chunk_size:
                chunk = self._create_document_chunk(
                    text=text,
                    source_file=source_file,
                    chunk_index=len(chunks),
                    page_data=page_data
                )
                chunks.append(chunk)
            else:
                # Create overlapping chunks
                start = 0
                while start < len(words):
                    end = min(start + chunk_size, len(words))
                    chunk_words = words[start:end]
                    chunk_text = ' '.join(chunk_words)
                    
                    chunk = self._create_document_chunk(
                        text=chunk_text,
                        source_file=source_file,
                        chunk_index=len(chunks),
                        page_data=page_data
                    )
                    chunks.append(chunk)
                    
                    # Move start position with overlap
                    if end >= len(words):
                        break
                    start = end - overlap
        
        logger.info(f"📦 Created {len(chunks)} intelligent chunks")
        return chunks
    
    def _create_document_chunk(self, text: str, source_file: str, 
                              chunk_index: int, page_data: Dict[str, Any]) -> DocumentChunk:
        """
        Create a DocumentChunk with comprehensive metadata
        
        Args:
            text: Chunk text
            source_file: Source file name
            chunk_index: Index of the chunk
            page_data: Page/slide metadata
        
        Returns:
            DocumentChunk object with empty embedding (to be filled later)
        """
        return DocumentChunk(
            id=str(uuid.uuid4()),
            text=text,
            embedding=None,  # Will be filled by embedding service
            metadata={
                'topic': self._extract_topic(text),
                'word_count': len(text.split()),
                'char_count': len(text),
                'page_type': page_data.get('type', 'unknown'),
                'original_page_number': page_data.get('page_number', page_data.get('slide_number', 0)),
                'processing_timestamp': datetime.now().isoformat()
            },
            source_file=source_file,
            chunk_index=chunk_index,
            created_at=datetime.now().isoformat()
        )
    
    def _extract_topic(self, text: str) -> str:
        """
        Extract topic/theme from text chunk
        
        Args:
            text: Text content
        
        Returns:
            Extracted topic
        """
        # Try to find a meaningful first sentence
        sentences = text.split('.')
        for sentence in sentences:
            sentence = sentence.strip()
            if 10 < len(sentence) < 80:
                return sentence
        
        # Fallback to first 60 characters
        return text[:60] + "..." if len(text) > 60 else text
    
    def process_and_store_document(self, file_path: str, 
                                  replace_existing: bool = True) -> Dict[str, Any]:
        """
        Complete RAG pipeline: Parse → Chunk → Embed → Store
        
        Args:
            file_path: Path to document file
            replace_existing: Whether to replace existing chunks from same file
        
        Returns:
            Processing results summary
        """
        try:
            filename = os.path.basename(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            if file_ext not in self.supported_formats:
                raise ValueError(f"Unsupported format: {file_ext}")
            
            logger.info(f"🚀 Processing document: {filename}")
            
            # Step 1: Extract text based on file type
            if file_ext == '.pdf':
                pages_data = self.extract_text_from_pdf(file_path)
            elif file_ext == '.pptx':
                pages_data = self.extract_text_from_ppt(file_path)
            else:
                raise ValueError(f"Unsupported format: {file_ext}")
            
            if not pages_data:
                raise ValueError("No text content found in document")
            
            # Step 2: Create intelligent chunks
            chunks = self.create_intelligent_chunks(pages_data, filename)
            
            # Step 3: Generate embeddings for chunks
            logger.info("🧠 Generating embeddings...")
            for i, chunk in enumerate(chunks):
                embedding_result = self.embedding_service.generate_embedding(chunk.text)
                chunk.embedding = embedding_result.embedding
                
                if (i + 1) % 10 == 0:
                    logger.info(f"Generated {i + 1}/{len(chunks)} embeddings")
            
            # Step 4: Remove existing chunks if requested
            if replace_existing:
                deleted_count = self.vector_db.delete_by_source(filename)
                if deleted_count > 0:
                    logger.info(f"🗑️ Removed {deleted_count} existing chunks")
            
            # Step 5: Store in vector database
            logger.info("💾 Storing in vector database...")
            chunk_ids = self.vector_db.add_document_chunks(chunks)
            
            # Prepare summary
            summary = {
                'success': True,
                'source_file': filename,
                'total_pages': len(pages_data),
                'total_chunks': len(chunks),
                'stored_chunk_ids': chunk_ids,
                'processing_timestamp': datetime.now().isoformat(),
                'embedding_dimension': self.embedding_service.embedding_dimension,
                'file_type': file_ext[1:]  # Remove the dot
            }
            
            logger.info(f"✅ Successfully processed {filename}: {len(chunks)} chunks stored")
            return summary
            
        except Exception as e:
            logger.error(f"❌ Failed to process document: {e}")
            return {
                'success': False,
                'error': str(e),
                'source_file': os.path.basename(file_path) if file_path else 'unknown'
            }
    
    def search_document_content(self, query: str, 
                               top_k: int = 5,
                               source_file_filter: Optional[str] = None) -> List[Dict[str, Any]]:
        """
        Search document content using semantic similarity
        
        Args:
            query: Search query
            top_k: Number of results to return
            source_file_filter: Optional filter by source file
        
        Returns:
            List of search results with metadata
        """
        try:
            # Generate query embedding
            query_embedding_result = self.embedding_service.generate_query_embedding(query)
            
            # Search vector database
            search_results = self.vector_db.similarity_search(
                query_embedding_result.embedding,
                top_k=top_k,
                source_file_filter=source_file_filter
            )
            
            # Format results
            formatted_results = []
            for result in search_results:
                formatted_results.append({
                    'text': result.chunk.text,
                    'similarity_score': result.similarity_score,
                    'source_file': result.chunk.source_file,
                    'chunk_index': result.chunk.chunk_index,
                    'metadata': result.chunk.metadata,
                    'chunk_id': result.chunk.id
                })
            
            logger.info(f"🔍 Found {len(formatted_results)} results for query: {query[:50]}...")
            return formatted_results
            
        except Exception as e:
            logger.error(f"❌ Search failed: {e}")
            raise
    
    def get_database_stats(self) -> Dict[str, Any]:
        """Get comprehensive database statistics"""
        try:
            return self.vector_db.get_collection_stats()
        except Exception as e:
            logger.error(f"❌ Failed to get database stats: {e}")
            return {'error': str(e)}


# Test function
def test_enhanced_parser():
    """Test the enhanced document parser"""
    try:
        print("🧪 Testing Enhanced Document Parser...")
        
        # Initialize parser (will use default services)
        parser = EnhancedDocumentParser()
        
        # Test with a sample text (simulating document processing)
        print("✅ Parser initialized successfully")
        
        # Get database stats
        stats = parser.get_database_stats()
        print(f"📊 Database stats: {stats}")
        
        print("🎉 Enhanced parser test passed!")
        
    except Exception as e:
        print(f"❌ Test failed: {e}")

if __name__ == "__main__":
    test_enhanced_parser()
